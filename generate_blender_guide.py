"""
Generate: Guide Animation 3D Blender — Référence Prompting Claude Code
Output: Guide_Animation_Blender_Claude.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Design tokens ───────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x2D, 0x2D, 0x44)
BG_CODE = RGBColor(0x23, 0x23, 0x3A)
BG_TABLE_ROW_A = RGBColor(0x22, 0x22, 0x38)
BG_TABLE_ROW_B = RGBColor(0x2A, 0x2A, 0x42)
BG_TABLE_HEADER = RGBColor(0x14, 0x14, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
RED_SOFT = RGBColor(0xF8, 0x71, 0x71)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=LIGHT_GRAY, bold=False, font_name=FONT_BODY,
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, lines, default_size=13):
    """lines: list of (text, font_size, color, bold, font_name, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else LIGHT_GRAY
        bold = line_data[3] if len(line_data) > 3 else False
        fname = line_data[4] if len(line_data) > 4 else FONT_BODY
        align = line_data[5] if len(line_data) > 5 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = fname
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_title(slide, text, subtitle=None):
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 text, font_size=28, color=WHITE, bold=True)
    # Green accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(1.0), Inches(2.5), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.15), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=MID_GRAY)


def add_code_block(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CODE
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x55)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)
    tf.margin_bottom = Inches(0.15)
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = GREEN
        p.font.name = FONT_CODE
        p.space_after = Pt(2)
    return shape


def add_table(slide, left, top, width, rows_data, col_widths, header=True):
    """rows_data: list of lists of strings. First row = header if header=True."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                         Inches(0.35 * n_rows))
    table = table_shape.table
    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            # Style
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = FONT_BODY
                if r_idx == 0 and header:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = GREEN
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.alignment = PP_ALIGN.LEFT
            # Cell fill
            if r_idx == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_HEADER
            elif r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_ROW_A
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_ROW_B
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return table_shape


def add_card(slide, left, top, width, height, title_text, body_lines):
    """Card with title + body lines inside a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = RGBColor(0x3F, 0x3F, 0x5C)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_right = Inches(0.15)
    tf.margin_bottom = Inches(0.1)
    # Title
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(13)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p.space_after = Pt(6)
    # Body
    for line in body_lines:
        p = tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1]
        else:
            p.text = line
            p.font.color.rgb = LIGHT_GRAY
        p.font.size = Pt(10.5)
        p.font.name = FONT_BODY
        p.space_after = Pt(2)
    return shape


def add_section_badge(slide, text, left=Inches(0.6), top=Inches(0.05)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(2.2), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x4A, 0xDE, 0x80)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_top = Pt(1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x0A, 0x0A, 0x1A)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


# ─── Build presentation ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
# Central title block
add_text_box(sl, Inches(1), Inches(1.8), Inches(11.3), Inches(1),
             "Guide Animation 3D Blender", font_size=42, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(2.9), Inches(11.3), Inches(0.6),
             "Reference Prompting pour Claude Code", font_size=22, color=GREEN,
             alignment=PP_ALIGN.CENTER)
# Decorative line
shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(5), Inches(3.7), Inches(3.3), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()
# Subtitle
add_text_box(sl, Inches(1), Inches(4.1), Inches(11.3), Inches(0.8),
             "Vocabulaire technique \u2022 Catalogue d'effets \u2022 Templates de prompts\n"
             "Tout ce qu'il faut pour piloter Blender via Claude Code",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
# Version
add_text_box(sl, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
             "v1.0 \u2014 Avril 2026", font_size=11, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Sommaire
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_title(sl, "Sommaire")

sections = [
    ("1", "Fondamentaux", "Concepts cl\u00e9s, keyframes, transformations, rep\u00e8res", GREEN),
    ("2", "Catalogue des effets", "Rotations, translations, \u00e9chelle, shape keys, cam\u00e9ra", BLUE),
    ("3", "Contr\u00f4le du mouvement", "Interpolation, F-Curves, contraintes, drivers", AMBER),
    ("4", "Prompting efficace", "Structure, lexique FR\u2192Blender, templates, anti-patterns", GREEN),
    ("5", "Export Web", "Bake, NLA, GLB, Draco, checklist", BLUE),
]
y = Inches(1.6)
for num, title, desc, color in sections:
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1.5), y, Inches(10.3), Inches(0.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = RGBColor(0x3F, 0x3F, 0x5C)
    shape.line.width = Pt(1)
    # Number badge
    badge = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(1.8), y + Inches(0.15), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(sl, Inches(2.6), y + Inches(0.08), Inches(8), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(sl, Inches(2.6), y + Inches(0.45), Inches(8), Inches(0.35),
                 desc, font_size=12, color=MID_GRAY)
    y += Inches(1.05)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Concepts clés
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 1 \u2014 FONDAMENTAUX")
add_title(sl, "Concepts cl\u00e9s de l'animation Blender")

concepts = [
    ("Keyframe", "Point cl\u00e9 dans la timeline qui enregistre\nla valeur d'une propri\u00e9t\u00e9 \u00e0 un instant donn\u00e9."),
    ("Timeline", "Ligne temporelle o\u00f9 sont plac\u00e9s les keyframes.\nUnit\u00e9 = frame (image)."),
    ("Frame Rate (FPS)", "Nombre d'images/seconde.\n24 fps (cin\u00e9ma), 30 fps (vid\u00e9o), 60 fps (jeu)."),
    ("Interpolation", "Calcul automatique des valeurs entre\n2 keyframes. Contr\u00f4l\u00e9 par les F-Curves."),
    ("F-Curve", "Courbe dans le Graph Editor qui d\u00e9finit\nl'\u00e9volution d'une valeur dans le temps."),
    ("Action", "Bloc r\u00e9utilisable d'animation.\nContient un ensemble de F-Curves."),
    ("NLA Editor", "Non-Linear Animation \u2014 empile et mixe\nplusieurs Actions comme des clips."),
]

x_positions = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.5)]
y_start = Inches(1.5)
for i, (term, desc) in enumerate(concepts):
    col = i % 4
    row = i // 4
    x = x_positions[col]
    y = y_start + row * Inches(2.5)
    add_card(sl, x, y, Inches(2.8), Inches(2.1), term, desc.split("\n"))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Les 3 transformations
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 1 \u2014 FONDAMENTAUX")
add_title(sl, "Les 3 transformations fondamentales")

transforms = [
    ("Location (G)", GREEN,
     ["D\u00e9placement de l'objet dans l'espace",
      "Axes : X (rouge), Y (vert), Z (bleu)",
      "Raccourci : G puis X/Y/Z",
      "Propri\u00e9t\u00e9 : object.location",
      ("Prompt : \"d\u00e9place le cube de 2m en Z\"", BLUE)]),
    ("Rotation (R)", AMBER,
     ["Rotation autour d'un axe",
      "Modes : Euler XYZ, Quaternion",
      "Raccourci : R puis X/Y/Z",
      "Propri\u00e9t\u00e9 : object.rotation_euler",
      ("Prompt : \"tourne le cube de 90\u00b0 en Z\"", BLUE)]),
    ("Scale (S)", BLUE,
     ["Mise \u00e0 l'\u00e9chelle (uniforme ou par axe)",
      "Uniform : m\u00eame valeur X=Y=Z",
      "Raccourci : S puis X/Y/Z",
      "Propri\u00e9t\u00e9 : object.scale",
      ("Prompt : \"double la taille en X\"", BLUE)]),
]

for i, (title, color, lines) in enumerate(transforms):
    x = Inches(0.5) + i * Inches(4.2)
    card = add_card(sl, x, Inches(1.5), Inches(3.9), Inches(3.2), title, lines)
    # Update title color per card
    card.text_frame.paragraphs[0].font.color.rgb = color

# Space concepts
add_card(sl, Inches(0.5), Inches(5.0), Inches(6), Inches(2.0),
         "Espace Local vs World",
         ["Local : axes relatifs \u00e0 l'orientation de l'objet",
          "World (Global) : axes fixes de la sc\u00e8ne",
          "Raccourci : appuyer 2x sur l'axe pour basculer (ex: G Z Z = Z local)",
          ("Prompt : \"d\u00e9place en Z local\" vs \"d\u00e9place en Z world\"", BLUE)])

add_card(sl, Inches(6.8), Inches(5.0), Inches(6), Inches(2.0),
         "Pivot Point",
         ["3D Cursor : pivot au curseur 3D",
          "Median Point : centre de la s\u00e9lection",
          "Individual Origins : chaque objet tourne sur lui-m\u00eame",
          "Active Element : pivot sur l'\u00e9l\u00e9ment actif",
          ("Prompt : \"tourne autour du curseur 3D\"", BLUE)])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Rotations
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 2 \u2014 CATALOGUE DES EFFETS")
add_title(sl, "Rotations", "Effets de rotation et pivotement")

rotation_data = [
    ["Effet", "Description", "Propri\u00e9t\u00e9 Blender", "Prompt FR"],
    ["Turntable", "Rotation continue autour de Z (360\u00b0 en boucle)", "rotation_euler.z + Cycles modifier", "\"tourne sur lui-m\u00eame\""],
    ["Spin rapide", "Rotation rapide sur X, Y ou Z", "rotation_euler.[x/y/z] multi-tours", "\"tourne rapidement en X\""],
    ["Orbit", "Objet tourne autour d'un autre (parent Empty)", "Empty parent + rotation_euler.z", "\"orbite autour du cube\""],
    ["Tilt / Inclinaison", "Rotation partielle (ex: 15\u00b0) sur X ou Y", "rotation_euler.x (valeur partielle)", "\"s'incline l\u00e9g\u00e8rement\""],
    ["Roll", "Rotation sur l'axe longitudinal (avant-arri\u00e8re)", "rotation_euler.y", "\"roule sur lui-m\u00eame\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), rotation_data,
          [Inches(1.8), Inches(3.8), Inches(3.5), Inches(3.2)])

add_code_block(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0),
               "# Exemples de prompts pour rotations\n"
               "\"Fais tourner le cube en continu autour de Z, en boucle sur 120 frames\"\n"
               "\"Le logo orbite autour du centre \u00e0 2m de distance, 1 tour complet en 3 secondes\"\n"
               "\"Incline la bo\u00eete de 15\u00b0 en X avec un ease in-out\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Translations
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 2 \u2014 CATALOGUE DES EFFETS")
add_title(sl, "Translations / D\u00e9placements", "Mouvements de position dans l'espace")

trans_data = [
    ["Effet", "Description", "Impl\u00e9mentation Blender", "Prompt FR"],
    ["Slide / Glissement", "Translation lin\u00e9aire sur un axe", "location.[x/y/z] lin\u00e9aire", "\"glisse vers la droite\""],
    ["Float / L\u00e9vitation", "Oscillation sinuso\u00efdale en Z", "location.z + Noise/Cycles modifier", "\"flotte doucement\""],
    ["Bounce / Rebond", "Rebond au sol (Y) avec d\u00e9c\u00e9l\u00e9ration", "location.z + Bounce easing", "\"rebondit au sol\""],
    ["Follow Path", "Suit une courbe B\u00e9zier", "Follow Path constraint", "\"suit le chemin\""],
    ["Throw / Lancer", "Trajectoire parabolique", "location.x lin\u00e9aire + location.z courbe", "\"lanc\u00e9 en arc de cercle\""],
    ["Zigzag", "Altern. gauche/droite en avan\u00e7ant", "location.x oscille + location.y lin\u00e9aire", "\"avance en zigzag\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), trans_data,
          [Inches(1.8), Inches(3.2), Inches(3.8), Inches(3.5)])

add_code_block(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7),
               "# Exemples de prompts pour translations\n"
               "\"Fais flotter la sph\u00e8re \u00e0 0.3m d'amplitude en Z, boucle douce sur 60 frames\"\n"
               "\"Le cube rebondit 3 fois puis s'arr\u00eate, hauteur initiale 2m\"\n"
               "\"D\u00e9place l'objet le long de la courbe B\u00e9zier en 4 secondes\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Scale
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 2 \u2014 CATALOGUE DES EFFETS")
add_title(sl, "\u00c9chelle / Scale", "Effets de taille et de d\u00e9formation")

scale_data = [
    ["Effet", "Description", "Impl\u00e9mentation Blender", "Prompt FR"],
    ["Pulse / Respiration", "Scale oscillant uniforme (1.0 \u2192 1.1 \u2192 1.0)", "scale XYZ + Cycles modifier", "\"pulse doucement\""],
    ["Pop-in / Apparition", "Scale 0 \u2192 1 avec overshoot", "scale 0\u21921.1\u21921.0, Back easing", "\"appara\u00eet en pop\""],
    ["Squash & Stretch", "Aplatir puis \u00e9tirer (non-uniforme)", "scale.z down + scale.xy up, puis inverse", "\"s'\u00e9crase et rebondit\""],
    ["Grow / Croissance", "Scale progressif de petit \u00e0 grand", "scale 0.1\u21921.0 lin\u00e9aire ou ease", "\"grandit progressivement\""],
    ["Shrink / R\u00e9tr\u00e9cir", "Scale progressif de grand \u00e0 petit", "scale 1.0\u21920.0", "\"r\u00e9tr\u00e9cit jusqu'\u00e0 dispara\u00eetre\""],
    ["Heartbeat", "Double pulse rapide + pause", "scale pulse 2x rapide + hold", "\"bat comme un c\u0153ur\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), scale_data,
          [Inches(1.8), Inches(3.2), Inches(3.8), Inches(3.5)])

add_code_block(sl, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5),
               "# Exemples de prompts pour scale\n"
               "\"Le logo pulse de 1.0 \u00e0 1.15 en boucle, rythme lent (90 frames par cycle)\"\n"
               "\"Pop-in du bouton : scale 0 \u2192 1.2 \u2192 1.0 sur 20 frames, easing Back\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Shape Keys / Morphing
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 2 \u2014 CATALOGUE DES EFFETS")
add_title(sl, "Shape Keys / Morphing", "D\u00e9formations de maillage")

add_card(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.8),
         "Principe",
         ["Les Shape Keys d\u00e9finissent des formes alternatives du maillage",
          "Basis = forme de r\u00e9f\u00e9rence (repos)",
          "Chaque key a une valeur (0.0 \u2192 1.0)",
          "Animer la valeur = morpher entre les formes",
          "Propri\u00e9t\u00e9 : object.data.shape_keys.key_blocks['Key'].value",
          ("Id\u00e9al pour : visages, d\u00e9formations organiques, UI anim\u00e9es", BLUE)])

add_card(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8),
         "Cas d'usage typiques",
         ["Expressions faciales (smile, blink, surprise)",
          "Transformation d'un objet en un autre",
          "Respiration / gonflement organique",
          "Boutons UI qui changent de forme",
          "Liquide qui se d\u00e9forme",
          ("Combinable avec scale + rotation pour plus d'impact", GREEN)])

add_code_block(sl, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5),
               "# Exemples de prompts pour Shape Keys\n"
               "\"Cr\u00e9e une shape key 'Smile' qui rel\u00e8ve les coins de la bouche, anime de 0 \u00e0 1 sur 30 frames\"\n"
               "\"Morphe la sph\u00e8re en cube via shape key, transition sur 2 secondes avec ease in-out\"\n"
               "\"Ajoute un shape key 'Breathe' qui gonfle le torse, boucle sinuso\u00efdale 0.0\u21920.6\u21920.0\"\n"
               "\"Anime les shape keys 'Blink_L' et 'Blink_R' en sync : 0\u21921\u21920 sur 8 frames, toutes les 3 secondes\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Combinaisons d'effets
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 2 \u2014 CATALOGUE DES EFFETS")
add_title(sl, "Combinaisons d'effets", "Mouvements compos\u00e9s et complexes")

combo_data = [
    ["Effet", "Composition", "Description", "Prompt FR"],
    ["Wobble", "Rotation noise XYZ", "Tremblement l\u00e9ger multi-axes", "\"tremble l\u00e9g\u00e8rement\""],
    ["Swing", "Rotation Z oscillante", "Balancement comme un pendule", "\"se balance doucement\""],
    ["Spiral / H\u00e9lice", "Rotation Z + Translation Z", "Monte ou descend en spirale", "\"monte en spirale\""],
    ["Pendulum", "Rotation avec pivot d\u00e9centr\u00e9", "Mouvement de pendule r\u00e9aliste", "\"oscille comme un pendule\""],
    ["Float + Rotate", "Location Z sinus + Rot Z", "L\u00e9vitation + rotation douce", "\"flotte et tourne lentement\""],
    ["Bounce + Squash", "Loc Z bounce + Scale squash", "Rebond r\u00e9aliste cartoon", "\"rebondit avec squash & stretch\""],
    ["Orbit + Scale", "Empty rotation + Scale pulse", "Orbite avec pulsation", "\"orbite en pulsant\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), combo_data,
          [Inches(1.6), Inches(2.8), Inches(3.8), Inches(4.1)])

add_code_block(sl, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.4),
               "# Exemples de prompts combin\u00e9s\n"
               "\"La mascotte flotte en Z (\u00b10.2m, 60 frames) tout en tournant lentement en Z (1 tour / 120 frames)\"\n"
               "\"Rebond cartoon : bounce en Z + squash au contact + stretch en l'air, 3 rebonds\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Mouvements de caméra
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 2 \u2014 CATALOGUE DES EFFETS")
add_title(sl, "Mouvements de cam\u00e9ra", "Effets cin\u00e9matographiques")

cam_data = [
    ["Mouvement", "Type", "Propri\u00e9t\u00e9", "Prompt FR"],
    ["Pan", "Rotation Y", "camera.rotation_euler.z", "\"panoramique horizontal\""],
    ["Tilt", "Rotation X", "camera.rotation_euler.x", "\"panoramique vertical\""],
    ["Dolly", "Translation Z (profondeur)", "camera.location.y (vers objet)", "\"la cam\u00e9ra s'approche\""],
    ["Truck", "Translation X (lat\u00e9ral)", "camera.location.x", "\"travelling lat\u00e9ral\""],
    ["Crane / Boom", "Translation Z (hauteur)", "camera.location.z", "\"mont\u00e9e en grue\""],
    ["Orbit cam", "Camera parent\u00e9e \u00e0 Empty", "Empty.rotation_euler.z", "\"cam\u00e9ra tourne autour\""],
    ["Dolly Zoom", "Dolly avant + FOV augmente", "location.y + lens", "\"effet Vertigo\""],
    ["Rack Focus", "Changement DoF focus", "dof.focus_distance", "\"mise au point sur l'arri\u00e8re-plan\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), cam_data,
          [Inches(1.6), Inches(2.8), Inches(3.5), Inches(4.4)])

add_code_block(sl, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.4),
               "# Exemples de prompts cam\u00e9ra\n"
               "\"Dolly avant lent vers le produit sur 90 frames, ease out \u00e0 la fin\"\n"
               "\"Orbit cam\u00e9ra 360\u00b0 autour de la sc\u00e8ne, hauteur 1.5m, 4 secondes\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Interpolation & Easing
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 3 \u2014 CONTR\u00d4LE DU MOUVEMENT")
add_title(sl, "Interpolation & Easing", "Contr\u00f4ler la vitesse et le ressenti du mouvement")

easing_data = [
    ["Type", "Comportement", "Cas d'usage", "Prompt FR"],
    ["Constant", "Pas de transition (saut instantan\u00e9)", "Switch on/off, visibilit\u00e9", "\"change instantan\u00e9ment\""],
    ["Linear", "Vitesse constante (robotique)", "Convoyeur, d\u00e9filement", "\"vitesse constante\""],
    ["Bezier", "Courbe lisse (d\u00e9faut Blender)", "Mouvement naturel g\u00e9n\u00e9ral", "\"mouvement fluide\""],
    ["Ease In", "D\u00e9marrage lent, acc\u00e9l\u00e8re", "D\u00e9part d'un objet lourd", "\"acc\u00e9l\u00e8re progressivement\""],
    ["Ease Out", "Rapide puis ralentit", "Arriv\u00e9e, atterrissage", "\"d\u00e9c\u00e9l\u00e8re en douceur\""],
    ["Ease In-Out", "Lent \u2192 rapide \u2192 lent", "La plupart des animations UI", "\"d\u00e9marrage et fin doux\""],
    ["Back", "D\u00e9passe la cible puis revient", "Pop-in, overshoot", "\"avec un l\u00e9ger d\u00e9passement\""],
    ["Elastic", "Rebond \u00e9lastique autour de la cible", "Cartoon, UI ludique", "\"rebond \u00e9lastique\""],
    ["Bounce", "Rebonds d\u00e9croissants", "Chute r\u00e9aliste, balle", "\"rebondit en arrivant\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), easing_data,
          [Inches(1.5), Inches(3.3), Inches(3.2), Inches(4.3)])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Modificateurs F-Curve
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 3 \u2014 CONTR\u00d4LE DU MOUVEMENT")
add_title(sl, "Modificateurs F-Curve", "Modifier le comportement des courbes d'animation")

cards_data = [
    ("Cycles", GREEN,
     ["Mode REPEAT : boucle identique \u00e0 l'infini",
      "Mode MIRROR : boucle en aller-retour",
      "Param\u00e8tres : Before/After, Count",
      ("Prompt : \"en boucle infinie\" / \"ping-pong\"", BLUE)]),
    ("Noise", AMBER,
     ["Ajoute une variation al\u00e9atoire \u00e0 la courbe",
      "Param\u00e8tres : Scale, Strength, Phase, Depth",
      "Strength faible = tremblement subtil",
      ("Prompt : \"avec du bruit\" / \"tremblement l\u00e9ger\"", BLUE)]),
    ("Stepped", RED_SOFT,
     ["Discr\u00e9tise le mouvement (paliers)",
      "Cr\u00e9e un effet stop-motion / saccad\u00e9",
      "Param\u00e8tres : Step Size, Offset",
      ("Prompt : \"style stop-motion\" / \"mouvement saccad\u00e9\"", BLUE)]),
    ("Limits", BLUE,
     ["Clampe les valeurs min et/ou max",
      "Emp\u00eache un objet de d\u00e9passer une zone",
      "Utile combin\u00e9 avec Noise ou Cycles",
      ("Prompt : \"limit\u00e9 entre 0 et 2m en Z\"", BLUE)]),
    ("Envelope", MID_GRAY,
     ["Module l'amplitude de la courbe",
      "Points de contr\u00f4le pour varier l'intensit\u00e9",
      "Fade in / fade out d'animation",
      ("Prompt : \"commence doucement, intensifie, puis calme\"", BLUE)]),
]

x_positions = [Inches(0.3), Inches(2.85), Inches(5.4), Inches(7.95), Inches(10.5)]
for i, (title, color, lines) in enumerate(cards_data):
    card = add_card(sl, x_positions[i], Inches(1.5), Inches(2.4), Inches(3.3), title, lines)
    card.text_frame.paragraphs[0].font.color.rgb = color

add_code_block(sl, Inches(0.3), Inches(5.1), Inches(12.5), Inches(2.1),
               "# Exemples de prompts F-Curve modifiers\n"
               "\"Rotation Z en boucle infinie, mode mirror (aller-retour)\"\n"
               "\"Ajoute un noise modifier sur location XYZ : strength 0.05, scale 3.0 (tremblement subtil)\"\n"
               "\"Animation stepped \u00e0 8 fps pour un rendu stop-motion\"\n"
               "\"Limite la position Z entre 0 et 5m, m\u00eame avec le bruit\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Contraintes (Constraints)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 3 \u2014 CONTR\u00d4LE DU MOUVEMENT")
add_title(sl, "Contraintes (Constraints)", "Automatiser les relations entre objets")

constr_data = [
    ["Constraint", "R\u00f4le", "Param\u00e8tres cl\u00e9s", "Prompt FR"],
    ["Track To", "L'objet regarde toujours la cible", "Target, Up axis, Track axis", "\"regarde toujours le cube\""],
    ["Follow Path", "L'objet suit une courbe", "Curve target, Offset, Follow", "\"suit la courbe B\u00e9zier\""],
    ["Copy Location", "Copie la position d'un autre objet", "Target, Axes, Influence", "\"suit la position de X\""],
    ["Copy Rotation", "Copie la rotation d'un autre objet", "Target, Axes, Mix mode", "\"copie la rotation de X\""],
    ["Copy Scale", "Copie l'\u00e9chelle d'un autre objet", "Target, Axes, Power", "\"m\u00eame taille que X\""],
    ["Limit Location", "Clampe la position (min/max)", "Min/Max XYZ, Owner space", "\"reste entre 0 et 3m en Z\""],
    ["Damped Track", "Orientation douce vers cible", "Target, Track axis", "\"suit du regard doucement\""],
    ["Child Of", "Parent dynamique (animable)", "Target, Influence (keyframable)", "\"attach\u00e9 \u00e0 la main, puis l\u00e2ch\u00e9\""],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), constr_data,
          [Inches(1.6), Inches(3.2), Inches(3.3), Inches(4.2)])

add_code_block(sl, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5),
               "# Exemples de prompts constraints\n"
               "\"La cam\u00e9ra Track To le personnage en permanence pendant l'orbite\"\n"
               "\"L'objet suit le path B\u00e9zier sur 120 frames, Follow Curve activ\u00e9\"\n"
               "\"Child Of : l'\u00e9p\u00e9e est attach\u00e9e \u00e0 la main frame 1\u219260, puis l\u00e2ch\u00e9e frame 61\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Drivers & Expressions
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 3 \u2014 CONTR\u00d4LE DU MOUVEMENT")
add_title(sl, "Drivers & Expressions", "Lier des propri\u00e9t\u00e9s entre elles par formule")

add_card(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5),
         "Principe des Drivers",
         ["Un driver remplace l'animation manuelle par une formule",
          "La propri\u00e9t\u00e9 driven est calcul\u00e9e automatiquement",
          "Variables : propri\u00e9t\u00e9s d'autres objets, frame, custom props",
          "Fonctions : sin(), cos(), min(), max(), clamp(), noise.random()",
          ("Propri\u00e9t\u00e9 driver\u00e9e = champ violet dans Blender", AMBER)])

add_card(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5),
         "Cas d'usage courants",
         ["Rotation pilot\u00e9e par la position (roue qui tourne en avan\u00e7ant)",
          "Scale proportionnel \u00e0 la distance",
          "Oscillation automatique via sin(frame / speed)",
          "Visibilit\u00e9 conditionnelle (si Z > seuil, visible)",
          ("Combinable avec constraints pour des rigs complexes", GREEN)])

driver_examples = [
    ["Formule", "Effet", "Prompt FR"],
    ["sin(frame / 30) * 0.5", "Oscillation sinuso\u00efdale (p\u00e9riode \u224860f)", "\"oscille en sinus\""],
    ["var_loc_z * 2.0", "Scale = 2x la hauteur Z", "\"grossit quand il monte\""],
    ["clamp(var, 0.0, 1.0)", "Valeur limit\u00e9e entre 0 et 1", "\"jamais en dessous de 0\""],
    ["noise.random() * 0.1", "Variation al\u00e9atoire \u00e0 chaque frame", "\"bruit al\u00e9atoire subtil\""],
    ["-var_rot_z", "Rotation inverse de la cible", "\"tourne \u00e0 l'oppos\u00e9 de X\""],
]
add_table(sl, Inches(0.5), Inches(4.3), Inches(12.3), driver_examples,
          [Inches(3.5), Inches(4.8), Inches(4.0)])

add_code_block(sl, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9),
               "# Prompt driver\n"
               "\"Ajoute un driver sur scale.z : sin(frame / 45) * 0.3 + 1.0 (pulsation automatique sans keyframes)\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Structure d'un prompt optimal
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 4 \u2014 PROMPTING EFFICACE")
add_title(sl, "Structure d'un prompt optimal", "Formuler des instructions claires pour Claude Code")

add_code_block(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
               "[OBJET]  +  [ACTION]  +  [AXE/DIRECTION]  +  [VITESSE/DUR\u00c9E]  +  [EASING]  +  [BOUCLE]")

add_text_box(sl, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.3),
             "Exemples de prompts bien structur\u00e9s :", font_size=14, color=WHITE, bold=True)

examples = [
    ("Basique", GREEN,
     "\"Fais tourner le cube autour de Z, 1 tour complet en 120 frames, en boucle\"",
     "Objet: cube | Action: rotation | Axe: Z | Dur\u00e9e: 120f | Boucle: oui"),
    ("Interm\u00e9diaire", BLUE,
     "\"Le logo flotte en Z (\u00b10.2m, 60 frames, ease in-out) et tourne lentement en Y (1 tour / 180 frames)\"",
     "2 actions | Amplitudes pr\u00e9cises | Timings diff\u00e9rents | Easing sp\u00e9cifi\u00e9"),
    ("Avanc\u00e9", AMBER,
     "\"Pop-in du produit : scale 0\u21921.15\u21921.0 sur 25 frames (Back easing), puis l\u00e9vitation Z \u00b10.1m en boucle, "
     "cam\u00e9ra dolly avant lent sur 90 frames\"",
     "3 animations s\u00e9quenc\u00e9es | Valeurs exactes | Easing | Boucle partielle"),
]

y = Inches(3.1)
for label, color, prompt, analysis in examples:
    badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), y, Inches(1.3), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_top = Pt(1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(sl, Inches(2.0), y - Inches(0.02), Inches(10.8), Inches(0.4),
                 prompt, font_size=11, color=WHITE, font_name=FONT_CODE)
    add_text_box(sl, Inches(2.0), y + Inches(0.38), Inches(10.8), Inches(0.3),
                 analysis, font_size=10, color=MID_GRAY)
    y += Inches(1.15)

# Tip box
tip = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8))
tip.fill.solid()
tip.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x2E)
tip.line.color.rgb = GREEN
tip.line.width = Pt(1)
tf = tip.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "ASTUCE : Plus le prompt est pr\u00e9cis (valeurs num\u00e9riques, axes, timing), plus le r\u00e9sultat sera conforme. " \
         "\u00c9vitez les descriptions vagues comme \"anime l'objet\" ou \"fais bouger le truc\"."
p.font.size = Pt(11)
p.font.color.rgb = GREEN
p.font.name = FONT_BODY

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Lexique FR → Terme Blender
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 4 \u2014 PROMPTING EFFICACE")
add_title(sl, "Lexique FR \u2192 Blender", "Correspondance termes fran\u00e7ais \u2192 propri\u00e9t\u00e9s Blender")

lexique_left = [
    ["Terme FR", "Propri\u00e9t\u00e9 / Technique Blender"],
    ["Faire tourner", "rotation_euler + Cycles modifier"],
    ["Flotter", "location.z sinuso\u00efdal + Cycles"],
    ["Rebondir", "location.z + Bounce easing"],
    ["Pulser", "scale XYZ oscillant + Cycles"],
    ["Trembler", "Noise modifier (faible Strength)"],
    ["Glisser", "location.[axe] lin\u00e9aire"],
    ["Se balancer", "rotation.z oscillant (pivot d\u00e9centr\u00e9)"],
    ["Appara\u00eetre", "scale 0\u21921 + Back easing (pop-in)"],
    ["Dispara\u00eetre", "scale 1\u21920 + Ease In"],
    ["Grandir", "scale progressif 0.1\u21921.0"],
    ["R\u00e9tr\u00e9cir", "scale progressif 1.0\u21920.0"],
    ["S'\u00e9craser", "Squash: scale.z down, scale.xy up"],
    ["S'\u00e9tirer", "Stretch: scale.z up, scale.xy down"],
    ["Orbiter", "Empty parent + rotation Z"],
    ["Suivre un chemin", "Follow Path constraint + courbe"],
    ["Regarder", "Track To constraint"],
    ["Monter en spirale", "rotation Z + location Z lin\u00e9aire"],
    ["Zig-zaguer", "location.x oscille + location.y avance"],
    ["Clignoter", "visibility keyframes (on/off)"],
]

lexique_right = [
    ["Terme FR", "Propri\u00e9t\u00e9 / Technique Blender"],
    ["Se d\u00e9former", "Shape Keys (morph 0\u21921)"],
    ["Changer de forme", "Shape Keys interpolation"],
    ["Respirer (objet)", "scale pulse lent + Cycles"],
    ["Battre (c\u0153ur)", "scale double-pulse + hold"],
    ["Lancer / jeter", "loc.x lin\u00e9aire + loc.z parabolique"],
    ["Acc\u00e9l\u00e9rer", "Ease In interpolation"],
    ["Ralentir", "Ease Out interpolation"],
    ["Rebond \u00e9lastique", "Elastic easing"],
    ["Stop-motion", "Stepped modifier (8-12 fps)"],
    ["En boucle", "Cycles modifier (REPEAT)"],
    ["Aller-retour", "Cycles modifier (MIRROR)"],
    ["Panoramique", "Camera rotation Y/Z"],
    ["Travelling", "Camera location X/Y"],
    ["Zoom avant", "Camera dolly (loc.y) ou lens"],
    ["Plongeante", "Camera en hauteur, tilt vers bas"],
    ["Contre-plongeante", "Camera basse, tilt vers haut"],
    ["Flou dynamique", "DOF focus_distance anim\u00e9"],
    ["Effet Vertigo", "Dolly + FOV inverse simult."],
    ["S'attacher / se d\u00e9tacher", "Child Of constraint (influence 0/1)"],
]

add_table(sl, Inches(0.3), Inches(1.5), Inches(6.2), lexique_left,
          [Inches(2.2), Inches(4.0)])
add_table(sl, Inches(6.7), Inches(1.5), Inches(6.2), lexique_right,
          [Inches(2.5), Inches(3.7)])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Templates de prompts (partie 1)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 4 \u2014 PROMPTING EFFICACE")
add_title(sl, "Templates de prompts (1/2)", "Prompts pr\u00eats \u00e0 l'emploi \u2014 basiques \u00e0 interm\u00e9diaires")

templates_1 = [
    ("BASIQUES", GREEN, [
        ("Turntable",
         "Fais tourner [OBJET] autour de Z, 1 tour complet en [N] frames, en boucle infinie"),
        ("Flottement",
         "Fais flotter [OBJET] en Z : amplitude \u00b1[X]m, cycle de [N] frames, ease in-out, boucle"),
        ("Bounce",
         "[OBJET] rebondit en Z depuis [H]m, 3 rebonds d\u00e9croissants, easing Bounce, dur\u00e9e [N] frames"),
        ("Pop-in",
         "[OBJET] appara\u00eet : scale 0 \u2192 1.15 \u2192 1.0 sur [N] frames, easing Back"),
        ("Pulse",
         "[OBJET] pulse : scale 1.0 \u2192 [X] \u2192 1.0, cycle [N] frames, ease in-out, boucle"),
    ]),
    ("INTERM\u00c9DIAIRES", BLUE, [
        ("Float + Rotate",
         "[OBJET] flotte en Z (\u00b1[X]m, [N]f) et tourne lentement en Y (1 tour / [M] frames)"),
        ("Orbit + Scale",
         "[OBJET] orbite autour de [CIBLE] \u00e0 [D]m, 1 tour en [N]f, pulse scale 1.0\u2192[X] en boucle"),
        ("Follow Path + Tilt",
         "[OBJET] suit [COURBE] en [N] frames, s'incline de [X]\u00b0 dans les virages"),
        ("S\u00e9quence entr\u00e9e",
         "[OBJET] slide depuis la gauche (loc.x -5\u21920, 30f, ease out) puis pop scale (0.8\u21921.1\u21921.0, 15f)"),
        ("Swing",
         "[OBJET] se balance en Z : rotation \u00b1[X]\u00b0, pivot d\u00e9centr\u00e9 en Y+[D], cycle [N]f, boucle mirror"),
    ]),
]

y = Inches(1.5)
for section_title, color, prompts in templates_1:
    add_text_box(sl, Inches(0.5), y, Inches(3), Inches(0.35),
                 section_title, font_size=13, color=color, bold=True)
    y += Inches(0.4)
    for name, template in prompts:
        add_text_box(sl, Inches(0.5), y, Inches(1.5), Inches(0.3),
                     name, font_size=10, color=color, bold=True)
        add_text_box(sl, Inches(2.1), y, Inches(10.5), Inches(0.3),
                     template, font_size=10, color=LIGHT_GRAY, font_name=FONT_CODE)
        y += Inches(0.38)
    y += Inches(0.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Templates de prompts (partie 2)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 4 \u2014 PROMPTING EFFICACE")
add_title(sl, "Templates de prompts (2/2)", "Prompts avanc\u00e9s et web-ready")

templates_2 = [
    ("AVANC\u00c9S", AMBER, [
        ("Dolly Zoom",
         "Cam\u00e9ra dolly avant (loc.y -10\u2192-3, 60f) + augmente focal de 35mm\u219270mm simult. (effet Vertigo)"),
        ("Squash & Stretch",
         "[OBJET] tombe en Z (2m\u21920, 15f ease in), squash au sol (scale .z=0.5 .xy=1.4, 3f), "
         "stretch rebond (.z=1.3 .xy=0.8, 5f), retour (1,1,1, 8f)"),
        ("Character Idle",
         "[PERSONNAGE] : location.z flotte \u00b10.05m (120f), rotation.z \u00b12\u00b0 (90f), "
         "scale pulse 1.0\u21921.02 (60f), le tout en boucle avec des p\u00e9riodes diff\u00e9rentes"),
        ("Camera Reveal",
         "Cam\u00e9ra commence derri\u00e8re l'objet (Y+5, Z+0.5), orbite 180\u00b0 en 120f (ease in-out), "
         "puis dolly avant l\u00e9ger (20f) pour reveal face"),
        ("Pendulum",
         "[OBJET] oscille comme un pendule : rotation X \u00b145\u00b0, pivot en Y+2m (origin d\u00e9centr\u00e9), "
         "cycle 90f, amortissement progressif via envelope modifier"),
    ]),
    ("WEB-READY (GLB)", GREEN, [
        ("Boucle parfaite",
         "[OBJET] rotation Z 360\u00b0 en [N] frames, frame 1 = frame N+1, Cycles modifier, bake, export GLB"),
        ("Hero product",
         "[PRODUIT] : turntable Z 360\u00b0 (120f) + l\u00e9vitation Z \u00b10.1m (60f) + scale pulse 1.0\u21921.05 (80f), "
         "boucle parfaite, bake all, export GLB Draco"),
        ("Logo anim\u00e9",
         "[LOGO] pop-in (scale 0\u21921.0, 20f Back) puis idle float+rotate en boucle, "
         "total 150f, bake, GLB avec Draco compression"),
        ("Loading spinner",
         "[OBJET] rotation Z continue 360\u00b0/60f, scale pulse 0.9\u21921.1 (30f mirror), "
         "parfaitement loopable, export GLB l\u00e9ger (< 500KB)"),
        ("Interactive hover",
         "[OBJET] idle: float Z \u00b10.05m (90f) + rot Y \u00b15\u00b0 (120f), "
         "hover state en Action s\u00e9par\u00e9e: scale 1.0\u21921.1 (15f ease out), NLA push, export GLB"),
    ]),
]

y = Inches(1.5)
for section_title, color, prompts in templates_2:
    add_text_box(sl, Inches(0.5), y, Inches(3), Inches(0.35),
                 section_title, font_size=13, color=color, bold=True)
    y += Inches(0.4)
    for name, template in prompts:
        add_text_box(sl, Inches(0.5), y, Inches(1.5), Inches(0.3),
                     name, font_size=10, color=color, bold=True)
        add_text_box(sl, Inches(2.1), y, Inches(10.5), Inches(0.35),
                     template, font_size=9.5, color=LIGHT_GRAY, font_name=FONT_CODE)
        y += Inches(0.42)
    y += Inches(0.15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Anti-patterns
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 4 \u2014 PROMPTING EFFICACE")
add_title(sl, "Anti-patterns", "\u00c0 \u00e9viter dans vos prompts d'animation")

antipatterns = [
    ("\u274c  \"Anime l'objet\"",
     "Trop vague \u2014 quel type d'animation ? rotation, translation, scale ?",
     "\u2705  \"Fais tourner le cube en Z, 1 tour en 120 frames, boucle\""),
    ("\u274c  \"Fais bouger le truc\"",
     "Pas d'objet pr\u00e9cis, pas de direction, pas de param\u00e8tres",
     "\u2705  \"D\u00e9place la sph\u00e8re de 3m en X sur 60 frames, ease out\""),
    ("\u274c  Oublier l'axe",
     "\"Fais tourner\" \u2014 autour de quel axe ? R\u00e9sultat impr\u00e9visible",
     "\u2705  Toujours sp\u00e9cifier X, Y ou Z (et local vs world si ambigu)"),
    ("\u274c  Pas de timing",
     "\"Fais rebondir\" \u2014 combien de frames ? quelle hauteur ?",
     "\u2705  \"Rebondit en Z depuis 2m, 3 rebonds, 90 frames total\""),
    ("\u274c  Oublier le bake/export",
     "L'animation fonctionne dans Blender mais ne s'exporte pas",
     "\u2705  \"Bake all actions, puis exporte en GLB avec Draco\""),
    ("\u274c  Animations > 5s pour le web",
     "Fichiers GLB trop lourds, temps de chargement excessif",
     "\u2705  Viser 2-4 secondes max, boucle parfaite, compression Draco"),
    ("\u274c  Trop d'effets simultan\u00e9s",
     "Rotation + translation + scale + morph = chaos ill\u00e9sible",
     "\u2705  Max 2-3 effets combin\u00e9s, avec des rythmes diff\u00e9rents"),
]

y = Inches(1.5)
for bad, why, good in antipatterns:
    row_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.4), y, Inches(12.5), Inches(0.75))
    row_shape.fill.solid()
    row_shape.fill.fore_color.rgb = BG_CARD
    row_shape.line.color.rgb = RGBColor(0x3F, 0x3F, 0x5C)
    row_shape.line.width = Pt(0.5)

    add_text_box(sl, Inches(0.6), y + Inches(0.05), Inches(3.5), Inches(0.3),
                 bad, font_size=11, color=RED_SOFT, bold=True)
    add_text_box(sl, Inches(4.2), y + Inches(0.05), Inches(4.5), Inches(0.3),
                 why, font_size=10, color=MID_GRAY)
    add_text_box(sl, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.3),
                 good, font_size=10, color=GREEN, font_name=FONT_CODE)
    y += Inches(0.82)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Préparer pour le web (Export)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 5 \u2014 EXPORT WEB")
add_title(sl, "Pr\u00e9parer pour le web", "Pipeline d'export animation \u2192 GLB")

steps = [
    ("1", "Bake All Actions", GREEN,
     ["Menu : Object \u2192 Animation \u2192 Bake Action",
      "Convertit drivers/constraints/modifiers en keyframes",
      "N\u00e9cessaire pour que l'animation s'exporte correctement",
      ("Prompt : \"bake toutes les animations\"", BLUE)]),
    ("2", "NLA Editor", BLUE,
     ["Push Down chaque Action dans le NLA",
      "Permet de g\u00e9rer plusieurs clips d'animation",
      "Utile pour s\u00e9parer idle / hover / click states",
      ("Prompt : \"push down les actions dans le NLA\"", BLUE)]),
    ("3", "Export GLB", AMBER,
     ["Format : glTF 2.0 (.glb = binaire compact)",
      "Options : Apply Modifiers \u2713, Bake Animation \u2713",
      "Sampling Rate : 1 (chaque frame) ou 2 (all\u00e9ger)",
      ("Prompt : \"exporte en GLB, apply modifiers\"", BLUE)]),
    ("4", "Compression Draco", GREEN,
     ["R\u00e9duit la taille du fichier de 50-80%",
      "Activer dans les options d'export GLB",
      "Compression level : 6 (bon compromis)",
      ("Prompt : \"active la compression Draco level 6\"", BLUE)]),
]

for i, (num, title, color, lines) in enumerate(steps):
    x = Inches(0.3) + i * Inches(3.2)
    # Number circle
    badge = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                x + Inches(0.05), Inches(1.5), Inches(0.45), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    card = add_card(sl, x, Inches(2.1), Inches(3.0), Inches(3.0), title, lines)
    card.text_frame.paragraphs[0].font.color.rgb = color

# Arrow connectors (visual)
for i in range(3):
    x = Inches(3.2) + i * Inches(3.2)
    add_text_box(sl, x, Inches(1.52), Inches(0.4), Inches(0.45),
                 "\u2192", font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom tips
add_code_block(sl, Inches(0.3), Inches(5.4), Inches(12.5), Inches(1.8),
               "# Prompt complet d'export\n"
               "\"Bake toutes les animations (Bake Action, Visual Keying activ\u00e9),\n"
               " push down dans le NLA,\n"
               " puis exporte en GLB (.glb) avec :\n"
               "   - Apply Modifiers activ\u00e9\n"
               "   - Bake Animation activ\u00e9, sampling rate 1\n"
               "   - Compression Draco level 6\n"
               "   - Fichier : produit_anime.glb\"")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Checklist finale
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_section_badge(sl, "SECTION 5 \u2014 EXPORT WEB")
add_title(sl, "Checklist finale", "R\u00e9sum\u00e9 des bonnes pratiques")

checklist_left = [
    ("\u2705", "Sp\u00e9cifier l'objet par son nom exact", GREEN),
    ("\u2705", "Pr\u00e9ciser le type d'animation (rotation, translation, scale...)", GREEN),
    ("\u2705", "Indiquer l'axe (X, Y, Z) et l'espace (local/world)", GREEN),
    ("\u2705", "Donner des valeurs num\u00e9riques (degr\u00e9s, m\u00e8tres, frames)", GREEN),
    ("\u2705", "Sp\u00e9cifier l'easing (Bezier, Ease In-Out, Back, Bounce...)", GREEN),
    ("\u2705", "Pr\u00e9ciser si l'animation doit boucler", GREEN),
    ("\u2705", "Limiter \u00e0 2-3 effets combin\u00e9s maximum", GREEN),
    ("\u2705", "Viser 2-4 secondes pour les animations web", GREEN),
]

checklist_right = [
    ("\u2705", "V\u00e9rifier le frame range (d\u00e9but et fin)", GREEN),
    ("\u2705", "S'assurer que frame 1 = frame N+1 pour les boucles", GREEN),
    ("\u2705", "Bake les animations avant export", GREEN),
    ("\u2705", "Utiliser le NLA pour les animations multi-states", GREEN),
    ("\u2705", "Exporter en GLB (pas glTF s\u00e9par\u00e9) pour le web", GREEN),
    ("\u2705", "Activer la compression Draco", GREEN),
    ("\u2705", "Tester la boucle dans un viewer GLB", GREEN),
    ("\u2705", "V\u00e9rifier la taille du fichier (< 1MB id\u00e9al)", GREEN),
]

y = Inches(1.5)
for icon, text, color in checklist_left:
    add_text_box(sl, Inches(0.5), y, Inches(0.4), Inches(0.35),
                 icon, font_size=14, color=color)
    add_text_box(sl, Inches(1.0), y + Inches(0.02), Inches(5.5), Inches(0.3),
                 text, font_size=12, color=LIGHT_GRAY)
    y += Inches(0.42)

y = Inches(1.5)
for icon, text, color in checklist_right:
    add_text_box(sl, Inches(6.8), y, Inches(0.4), Inches(0.35),
                 icon, font_size=14, color=color)
    add_text_box(sl, Inches(7.3), y + Inches(0.02), Inches(5.5), Inches(0.3),
                 text, font_size=12, color=LIGHT_GRAY)
    y += Inches(0.42)

# Final tip
tip = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8))
tip.fill.solid()
tip.fill.fore_color.rgb = BG_CARD
tip.line.color.rgb = GREEN
tip.line.width = Pt(1.5)
tf = tip.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.25)
tf.margin_top = Inches(0.15)

p = tf.paragraphs[0]
p.text = "WORKFLOW R\u00c9SUM\u00c9"
p.font.size = Pt(14)
p.font.color.rgb = GREEN
p.font.bold = True
p.space_after = Pt(8)

steps_text = [
    "1. D\u00e9crire l'animation avec pr\u00e9cision (objet + action + axe + timing + easing + boucle)",
    "2. Claude Code g\u00e9n\u00e8re le script Python Blender",
    "3. V\u00e9rifier le r\u00e9sultat dans le viewport (jouer la timeline)",
    "4. Ajuster si n\u00e9cessaire (\"rends le plus lent\", \"ajoute un ease out\", ...)",
    "5. Bake \u2192 NLA \u2192 Export GLB Draco \u2192 Int\u00e9grer dans le site web",
]
for step in steps_text:
    p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = FONT_BODY
    p.space_after = Pt(3)

# ─── Save ─────────────────────────────────────────────────────────────────────
output_path = r"C:\00 - CLAUDE\TEST1\Guide_Animation_Blender_Claude.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
